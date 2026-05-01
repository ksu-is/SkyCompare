from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

RED      = RGBColor(0xC0, 0x39, 0x2B)
DARK_BG  = RGBColor(0x11, 0x11, 0x11)
BORDER   = RGBColor(0x33, 0x33, 0x33)
LIGHT    = RGBColor(0xF0, 0xF0, 0xF0)
GRAY     = RGBColor(0xAA, 0xAA, 0xAA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)


def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def tb(slide, x, y, w, h, text, size, color,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = wrap
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


# ── Image dimensions ─────────────────────────────────────────────────────
search_img  = "Skycompare2.png"   # portrait  (search page)
results_img = "Skycompare1.png"   # landscape (results page)

sw, sh = Image.open(search_img).size
rw, rh = Image.open(results_img).size

# Both screenshots at the same rendered height
TARGET_H = Inches(3.25)
search_W  = TARGET_H * sw / sh
results_W = TARGET_H * rw / rh

# ── Slide setup ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
rect(slide, 0, 0, prs.slide_width, prs.slide_height, DARK_BG)
rect(slide, 0, 0, prs.slide_width, Inches(0.07), RED)           # top bar
rect(slide, 0, Inches(7.43), prs.slide_width, Inches(0.07), RED) # bottom bar

# ── Left column ──────────────────────────────────────────────────────────
LEFT_W = Inches(4.7)

tb(slide, Inches(0.35), Inches(0.16), LEFT_W, Inches(1.0),
   "SkyCompare", 50, RED, bold=True)

tb(slide, Inches(0.35), Inches(1.18), LEFT_W, Inches(0.4),
   "Your Priorities. Your Flight.", 14, GRAY, italic=True)

tb(slide, Inches(0.35), Inches(1.72), LEFT_W, Inches(0.32),
   "Jaylen Early", 13, LIGHT, bold=True)

tb(slide, Inches(0.35), Inches(2.0), LEFT_W, Inches(0.25),
   "IS Project  ·  Spring 2026", 10, GRAY)

tb(slide, Inches(0.35), Inches(2.22), LEFT_W, Inches(0.25),
   "github.com/ksu-is/SkyCompare", 9, RGBColor(0xC0, 0x39, 0x2B))

rect(slide, Inches(0.35), Inches(2.52), Inches(4.4), Inches(0.03), BORDER)

features = [
    "Live prices from Google Flights via SerpAPI",
    "AviationStack for real departure & arrival times",
    "Airport autocomplete — 80+ US airports",
    "3-priority ranking system powers the SkyScore",
    "Direct booking links pre-filled with route & date",
    "Red & black dark UI built on Bootstrap 5",
]

feat_box = slide.shapes.add_textbox(Inches(0.35), Inches(2.66), Inches(4.55), Inches(4.0))
feat_box.text_frame.word_wrap = True
for i, feat in enumerate(features):
    p = feat_box.text_frame.paragraphs[0] if i == 0 else feat_box.text_frame.add_paragraph()
    p.space_before = Pt(6)
    bullet = p.add_run()
    bullet.text = "▸  "
    bullet.font.size = Pt(11)
    bullet.font.color.rgb = RED
    bullet.font.bold = True
    body = p.add_run()
    body.text = feat
    body.font.size = Pt(11)
    body.font.color.rgb = LIGHT

# Vertical divider
rect(slide, Inches(5.0), Inches(0.12), Inches(0.025), Inches(7.26), BORDER)

# ── Screenshots ──────────────────────────────────────────────────────────
SCREEN_TOP  = Inches(1.0)
GAP         = Inches(0.22)
SEARCH_X    = Inches(5.2)
RESULTS_X   = SEARCH_X + search_W + GAP

# Thin red label bars above each screenshot
rect(slide, SEARCH_X, SCREEN_TOP - Inches(0.28), search_W, Inches(0.25), RGBColor(0x2A,0x0A,0x0A), RED)
tb(slide, SEARCH_X, SCREEN_TOP - Inches(0.27), search_W, Inches(0.24),
   "SEARCH INTERFACE", 8, LIGHT, bold=True, align=PP_ALIGN.CENTER)

rect(slide, RESULTS_X, SCREEN_TOP - Inches(0.28), results_W, Inches(0.25), RGBColor(0x2A,0x0A,0x0A), RED)
tb(slide, RESULTS_X, SCREEN_TOP - Inches(0.27), results_W, Inches(0.24),
   "RANKED RESULTS", 8, LIGHT, bold=True, align=PP_ALIGN.CENTER)

# Actual screenshots
slide.shapes.add_picture(search_img,  SEARCH_X,  SCREEN_TOP, height=TARGET_H)
slide.shapes.add_picture(results_img, RESULTS_X, SCREEN_TOP, height=TARGET_H)

# Thin border around each screenshot (transparent fill so image shows through)
for bx, bw in [(SEARCH_X, search_W), (RESULTS_X, results_W)]:
    s = slide.shapes.add_shape(1, bx, SCREEN_TOP, bw, TARGET_H)
    s.fill.background()
    s.line.color.rgb = BORDER

# ── Bottom tagline ───────────────────────────────────────────────────────
BOTTOM_Y = SCREEN_TOP + TARGET_H + Inches(0.18)
tb(slide, Inches(5.2), BOTTOM_Y, prs.slide_width - Inches(5.4), Inches(0.9),
   "SkyCompare helps travelers cut through the noise — not just finding the cheapest "
   "flight, but the right flight based on what matters most to them.",
   10, GRAY, italic=True, wrap=True)

prs.save("SkyCompare_Slide.pptx")
print("Done — SkyCompare_Slide.pptx saved.")
